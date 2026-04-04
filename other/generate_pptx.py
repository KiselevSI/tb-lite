#!/usr/bin/env python3
"""Generate TB-Lite conference presentation as PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Colors ───
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BLUE  = RGBColor(0x1A, 0x3C, 0x6E)
ACCENT     = RGBColor(0x2E, 0x5E, 0xA8)
TEXT_DARK  = RGBColor(0x2D, 0x2D, 0x2D)
TEXT_GRAY  = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
GREEN_ACC  = RGBColor(0x27, 0x8B, 0x4A)
RED_ACC    = RGBColor(0xC0, 0x39, 0x2B)
DARK_RED   = RGBColor(0x8E, 0x24, 0x1E)
ORANGE_ACC = RGBColor(0xE6, 0x7E, 0x22)
BORDER     = RGBColor(0xCC, 0xCC, 0xCC)
BAR_BLUE   = RGBColor(0x3A, 0x7C, 0xBD)
TBL_HEAD   = RGBColor(0x2E, 0x5E, 0xA8)
TBL_ALT    = RGBColor(0xF5, 0xF8, 0xFC)

FONT = "Times New Roman"
TOTAL = 11

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ══════════════ HELPERS ══════════════
def _bg(sl):
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    b = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
    b.fill.solid(); b.fill.fore_color.rgb = DARK_BLUE; b.line.fill.background()


def _num(sl, n):
    tf = sl.shapes.add_textbox(Inches(12), Inches(7.1), Inches(1.2), Inches(0.3)).text_frame
    p = tf.paragraphs[0]; p.text = f"{n} / {TOTAL}"
    p.font.size = Pt(12); p.font.name = FONT; p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.RIGHT


def _title(sl, text, sz=Pt(32)):
    tf = sl.shapes.add_textbox(Inches(0.7), Inches(0.25), Inches(11.5), Inches(0.65)).text_frame
    p = tf.paragraphs[0]; p.text = text
    p.font.size = sz; p.font.name = FONT; p.font.color.rgb = DARK_BLUE; p.font.bold = True
    r = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.95), Inches(11.5), Inches(0.025))
    r.fill.solid(); r.fill.fore_color.rgb = ACCENT; r.line.fill.background()


def _tb(sl, l, t, w, h, text, sz=Pt(18), clr=TEXT_DARK, bold=False, italic=False,
        align=PP_ALIGN.LEFT):
    tf = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = sz; p.font.name = FONT; p.font.color.rgb = clr
    p.font.bold = bold; p.font.italic = italic; p.alignment = align
    return tf


def _bullets(sl, l, t, w, h, items, sz=Pt(17), clr=TEXT_DARK, sp=Pt(5)):
    tf = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True
    for i, txt in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u25B8  {txt}"
        p.font.size = sz; p.font.name = FONT; p.font.color.rgb = clr; p.space_after = sp
    return tf


def _nbullets(sl, l, t, w, h, items, sz=Pt(18), sp=Pt(10)):
    tf = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True
    for i, txt in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i+1}.  {txt}"
        p.font.size = sz; p.font.name = FONT; p.font.color.rgb = TEXT_DARK; p.space_after = sp
    return tf


def _card(sl, l, t, w, h, bc=BORDER):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = bc; s.line.width = Pt(1.5)
    return s


def _stat(sl, l, t, w, h, label, value, desc, vc=ACCENT, bc=ACCENT):
    _card(sl, l, t, w, h, bc)
    _tb(sl, l+.1, t+.06, w-.2, .2, label.upper(), Pt(10), LIGHT_GRAY, bold=True, align=PP_ALIGN.CENTER)
    _tb(sl, l+.1, t+.3, w-.2, .35, value, Pt(22), vc, bold=True, align=PP_ALIGN.CENTER)
    if desc:
        _tb(sl, l+.1, t+.65, w-.2, .2, desc, Pt(10), TEXT_GRAY, align=PP_ALIGN.CENTER)


def _bar(sl, l, t, max_w, frac, h, clr):
    bw = max(max_w * frac, 0.04)
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(bw), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()
    return bw


def _tbl(sl, l, t, w, data, cw=None, fs=Pt(13), rh=.35):
    R, C = len(data), len(data[0])
    sh = sl.shapes.add_table(R, C, Inches(l), Inches(t), Inches(w), Inches(R * rh))
    tbl = sh.table
    for i, row in enumerate(data):
        for j, v in enumerate(row):
            c = tbl.cell(i, j); c.text = str(v); c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = TBL_HEAD if i == 0 else (TBL_ALT if i % 2 == 0 else WHITE)
            for pg in c.text_frame.paragraphs:
                pg.font.size = fs; pg.font.name = FONT
                pg.font.color.rgb = WHITE if i == 0 else TEXT_DARK
                pg.font.bold = (i == 0)
                pg.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
    if cw:
        for j, w_ in enumerate(cw):
            tbl.columns[j].width = Inches(w_)
    return sh


def _hbox(sl, l, t, w, h, bc=ACCENT):
    """Highlight box with light blue fill."""
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = TBL_ALT
    s.line.color.rgb = bc; s.line.width = Pt(1.5)


def _pipe(sl, l, t, name, tools, num, bc=BORDER, nc=TEXT_DARK, w=1.4, h=0.85):
    _card(sl, l, t, w, h, bc)
    _tb(sl, l, t+.05, w, .15, str(num), Pt(9), LIGHT_GRAY, align=PP_ALIGN.CENTER)
    _tb(sl, l, t+.22, w, .22, name, Pt(12), nc, bold=True, align=PP_ALIGN.CENTER)
    _tb(sl, l, t+.46, w, .35, tools, Pt(8), TEXT_GRAY, align=PP_ALIGN.CENTER)


def _line(sl, x1, y1, x2, y2, clr=GREEN_ACC, w=Pt(2)):
    c = sl.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = clr; c.line.width = w


# ══════════════════════════════════════════════════════════════
#  SLIDE 1 — ТИТУЛЬНЫЙ
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)

_tb(sl, 1.0, 1.5, 11.3, 1.0,
    "TB-Lite: воспроизводимый Nextflow-пайплайн\n"
    "для анализа геномов Mycobacterium tuberculosis",
    Pt(36), DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)

_tb(sl, 1.0, 3.0, 11.3, .5,
    "С. И. Киселев\u00B9\u00B7\u00B2, Е. А. Шитиков\u00B2",
    Pt(24), TEXT_DARK, align=PP_ALIGN.CENTER)

tf = sl.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10.3), Inches(1.6)).text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("\u00B9 Московский физико-технический институт "
          "(национальный исследовательский университет)")
p.font.size = Pt(16); p.font.name = FONT; p.font.color.rgb = TEXT_GRAY
p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(8)
p2 = tf.add_paragraph()
p2.text = ("\u00B2 ФГБУ «Федеральный научно-клинический центр "
           "физико-химической медицины\nим. Ю. М. Лопухина» ФМБА России")
p2.font.size = Pt(16); p2.font.name = FONT; p2.font.color.rgb = TEXT_GRAY
p2.alignment = PP_ALIGN.CENTER

_tb(sl, 1.0, 6.2, 11.3, .4, "2025", Pt(20), LIGHT_GRAY, align=PP_ALIGN.CENTER)
_num(sl, 1)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2 — АКТУАЛЬНОСТЬ (только ТБ + WGS, без пайплайна)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _title(sl, "Актуальность")

# Левая колонка: проблема ТБ
_tb(sl, .7, 1.2, 5.5, .35, "Туберкулёз: масштаб проблемы", Pt(22), ACCENT, bold=True)

_bullets(sl, .7, 1.65, 5.8, 2.2, [
    "Инфекционное заболевание, вызываемое\nмикобактериями туберкулёзного комплекса (MTBC)",
    "7,5 млн новых случаев, 1,3 млн смертей\nежегодно (ВОЗ, 2024)",
    "Успешность лечения: 88% (DS-TB), 63% (MDR-TB)",
    "Длительность терапии устойчивых форм: до 24 мес.",
    "Распространение лекарственно-устойчивых форм\nосложняет контроль инфекции",
], Pt(17))

# Генетика MTBC
_tb(sl, .7, 4.2, 5.5, .35, "Генетические особенности MTBC", Pt(22), ACCENT, bold=True)

_bullets(sl, .7, 4.65, 5.8, 2.0, [
    "Геном ~4,4 м.п.о., GC-состав 65%",
    "Клональная структура популяции, нет ГПГ",
    "Частота мутаций: 0,3–0,5 SNP на геном в год",
    "Сходство нуклеотидных последовательностей >99,5%",
], Pt(17))

# Правая колонка: WGS
_tb(sl, 7.0, 1.2, 5.8, .35, "Полногеномное секвенирование — «всё в одном»", Pt(22), ACCENT, bold=True)

_bullets(sl, 7.0, 1.65, 5.8, 1.8, [
    "Молекулярное типирование (линии, сполиготипы)",
    "Маркеры лекарственной устойчивости",
    "Популяционно-геномные исследования",
    "Заменяет традиционные методы:\n  IS6110-RFLP, MIRU-VNTR, сполиготипирование, LSP",
], Pt(17))

# Блок «Потребность»
_hbox(sl, 7.0, 3.8, 5.8, 1.6, ACCENT)
_tb(sl, 7.15, 3.9, 5.5, .3, "Потребность", Pt(18), DARK_BLUE, bold=True)
_tb(sl, 7.15, 4.3, 5.5, 1.0,
    "Для рутинной работы с данными WGS по-прежнему\n"
    "востребованы простые и воспроизводимые\n"
    "биоинформатические конвейеры, объединяющие\n"
    "ключевые этапы анализа в одном запуске\n"
    "и выдающие стандартизованный результат",
    Pt(16), TEXT_GRAY)

_num(sl, 2)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3 — ЦЕЛЬ И ЗАДАЧИ
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _title(sl, "Цель и задачи")

_hbox(sl, .7, 1.2, 11.5, 1.2, DARK_BLUE)
_tb(sl, .85, 1.3, 11, .28, "Цель работы", Pt(22), DARK_BLUE, bold=True)
_tb(sl, .85, 1.65, 11, .65,
    "Разработка и апробация пайплайна TB-Lite для автоматизированной\n"
    "обработки данных полногеномного секвенирования Mycobacterium tuberculosis",
    Pt(20), TEXT_DARK, italic=True)

_tb(sl, .7, 2.7, 5, .35, "Задачи:", Pt(22), DARK_BLUE, bold=True)

_nbullets(sl, .7, 3.15, 11.5, 4.0, [
    "Реализовать Nextflow-пайплайн, объединяющий ключевые этапы\n"
    "анализа WGS данных M. tuberculosis: от контроля качества\n"
    "до генотипирования",
    "Интегрировать модули молекулярного типирования:\n"
    "филогенетические линии (TBLG), споликотипы (SpoTyping),\n"
    "лекарственная устойчивость (TB-Profiler), делеции (RD),\n"
    "IS6110 (ISMapper)",
    "Разработать модуль детекции смешанных инфекций (TB-Mix)\n"
    "и трёхуровневую систему фильтрации образцов",
    "Провести апробацию на коллекции >80 000 публичных\n"
    "геномов M. tuberculosis",
], Pt(18))

_num(sl, 3)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4 — ОБЗОР ПАЙПЛАЙНА (бывший слайд 1)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _title(sl, "TB-Lite: обзор пайплайна")

_tb(sl, .7, 1.05, 11, .3,
    "Nextflow-пайплайн для геномного анализа Mycobacterium tuberculosis",
    Pt(18), TEXT_GRAY, italic=True)

_bullets(sl, .7, 1.5, 5.8, 3.2, [
    "8 этапов, 26 процессов — модульная DSL2-архитектура",
    "Вход: FASTQ (paired-end / single-end)",
    "Референс: H37Rv (NC_000962.3)",
    "Трёхуровневая фильтрация\n  (покрытие, выравнивание, контаминация)",
    "Детекция смешанных инфекций (TB-Mix)",
    "Когортный анализ при > 1 образце",
], Pt(17))

_tb(sl, 7.0, 1.5, 5.5, .3, "Результаты анализа:", Pt(16), TEXT_GRAY, bold=True)
_bullets(sl, 7.0, 1.9, 5.5, 2.8, [
    "Филогенетическая линия (TBLG: L1–L7)",
    "Споликотип (SpoTyping)",
    "Лекарственная устойчивость (TB-Profiler)",
    "Регионы делеций RD (RD-анализ)",
    "Позиции IS6110 (ISMapper)",
    "Итоговые таблицы: XLSX, TSV",
], Pt(17))

# Теги внизу
tags = [("Nextflow DSL2", ACCENT), ("Docker", GREEN_ACC), ("Singularity", GREEN_ACC),
        ("Kubernetes", TEXT_GRAY), ("HPC-ready", ORANGE_ACC)]
tx = .7
for tag, tc in tags:
    tw = len(tag) * .11 + .4
    _card(sl, tx, 5.3, tw, .35, tc)
    _tb(sl, tx, 5.33, tw, .3, tag, Pt(11), tc, align=PP_ALIGN.CENTER)
    tx += tw + .15
_num(sl, 4)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5 — АРХИТЕКТУРА (бывший слайд 2)
#  Точные координаты из оригинальной HTML-диаграммы
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _title(sl, "Архитектура пайплайна")

# Параметры сетки (оригинальные)
Y  = 2.9       # Y верхнего края блоков основного ряда
G  = 1.55      # шаг по X между блоками
X0 = 0.3       # начальный X
BW = 1.4       # ширина блока
BH = 0.85      # высота блока
QC_Y = 1.4     # Y для блока QC (выше основного ряда)
AN_Y = 4.7     # Y для блока ANN_TABLE (ниже)

boxes = [
    ("FASTQ",      "входные данные",               "вход", BORDER,     TEXT_GRAY),
    ("TRIMMING",   "fastp",                         "1",    GREEN_ACC,  GREEN_ACC),
    ("MAPPING",    "BWA-mem, samtools\nPicard",      "3",    BORDER,     TEXT_DARK),
    ("FILTER",     "TB-Mix, Picard\nsamtools",       "4",    ORANGE_ACC, ORANGE_ACC),
    ("CALL VAR",   "Freebayes, SnpEff\nBCFtools",    "5",    BORDER,     TEXT_DARK),
    ("GENOTYPING", "TBLG, SpoTyping, Mosdepth\nTB-Profiler, RD, ISMapper",
                                                     "6",    BORDER,     TEXT_DARK),
    ("REPORTS",    "MultiQC, FINAL_TABLE\nTB_PLATFORM_TABLES",
                                                     "7",    GREEN_ACC,  GREEN_ACC),
]

for i, (nm, tl, nu, bc, nc) in enumerate(boxes):
    _pipe(sl, X0 + i * G, Y, nm, tl, nu, bc, nc)

# QC (над основным рядом)
_pipe(sl, X0 + 2 * G, QC_Y, "QC", "FastQC", "2", ACCENT, ACCENT)

# ANN_TABLE (под основным рядом)
_pipe(sl, X0 + 4 * G, AN_Y, "ANN_TABLE", "bcftools merge\nSnpEff, SnpSift", "8", ACCENT, ACCENT)

# ── Стрелки основного потока (горизонтальные, зелёные) ──
for i in range(6):
    x1 = X0 + i * G + BW           # правый край текущего блока
    x2 = X0 + (i + 1) * G          # левый край следующего
    y_mid = Y + BH / 2
    _line(sl, x1, y_mid, x2, y_mid, GREEN_ACC, Pt(2))

# ── TRIMMING → QC (диагональ вверх) ──
cx_trim = X0 + 1 * G + BW / 2      # центр X блока TRIMMING
cx_qc   = X0 + 2 * G + BW / 2      # центр X блока QC
_line(sl, cx_trim, Y, cx_qc, QC_Y + BH, ACCENT, Pt(1.5))

# ── QC → REPORTS (длинная линия через верх) ──
cx_rep = X0 + 6 * G + BW / 2       # центр X блока REPORTS
_line(sl, X0 + 2 * G + BW, QC_Y + BH / 2, cx_rep, Y, ACCENT, Pt(1))

# ── FILTER → GENOTYPING (дуга внизу) ──
cx_filt = X0 + 3 * G + BW / 2      # центр X блока FILTER
cx_geno = X0 + 5 * G + BW / 2      # центр X блока GENOTYPING
_line(sl, cx_filt, Y + BH, cx_geno, Y + BH, ORANGE_ACC, Pt(1.5))

# ── CALLVAR → ANN_TABLE (вертикаль вниз) ──
cx_cv = X0 + 4 * G + BW / 2        # центр X блока CALL VAR
_line(sl, cx_cv, Y + BH, cx_cv, AN_Y, ACCENT, Pt(1.5))

# ── Подписи на стрелках ──
_tb(sl, X0 + 3*G - .1, Y - .22, 1.5, .2, "bam_good", Pt(8), TEXT_GRAY, align=PP_ALIGN.CENTER)
_tb(sl, X0 + 4*G - .1, Y - .22, 1.5, .2, "vcf", Pt(8), TEXT_GRAY, align=PP_ALIGN.CENTER)
_tb(sl, X0 + 2*G + BW + .1, 1.5, 2, .2, "fastqc_reports", Pt(8), ACCENT)
_tb(sl, X0 + 3*G + .3, Y + BH + .03, 3, .2, "trimmed_good + bam_good", Pt(8), ORANGE_ACC)
_tb(sl, cx_cv + .1, 4.35, 1.5, .2, "> 1 sample", Pt(8), TEXT_GRAY)

# ── Легенда ──
for x, c, t in [(.3, GREEN_ACC, "— основной поток"),
                (3.0, ACCENT, "— дополнительные данные"),
                (6.5, ORANGE_ACC, "— фильтрованные данные")]:
    ln = sl.shapes.add_connector(1, Inches(x), Inches(6.35), Inches(x+.2), Inches(6.35))
    ln.line.color.rgb = c; ln.line.width = Pt(2)
    _tb(sl, x+.25, 6.23, 2.5, .2, t, Pt(9), TEXT_GRAY)

_num(sl, 5)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6 — КОНТРОЛЬ КАЧЕСТВА И ФИЛЬТРАЦИЯ (бывший слайд 3)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _title(sl, "Контроль качества и фильтрация")

# Поток: BAM → 3 модуля → Решение
flow = [
    ("BAM-файл", "Картированные риды\nBWA-mem + Picard MarkDups"),
    ("3 модуля параллельно", "Picard WGS\nSamtools stats\nTB-Mix"),
    ("Решение", "PASS → дальнейший анализ\nFAIL → bad_reads_low_coverage.txt"),
]
for i, (title, desc) in enumerate(flow):
    x = .5 + i * 4.0
    _card(sl, x, 1.2, 3.3, 1.1, BORDER)
    _tb(sl, x+.15, 1.25, 3.0, .25, title, Pt(14), ACCENT, bold=True)
    _tb(sl, x+.15, 1.55, 3.0, .7, desc, Pt(12), TEXT_DARK)

for i in range(2):
    _tb(sl, 3.5 + i*4.0, 1.5, .5, .3, "\u2192", Pt(24), GREEN_ACC, align=PP_ALIGN.CENTER)

# Пороги
_tb(sl, .7, 2.6, 10, .3, "Пороговые значения фильтрации", Pt(18), ACCENT, bold=True)

for i, (lbl, val, desc, clr) in enumerate([
    ("Медианное покрытие", "\u2265 30X", "Picard CollectWgsMetrics\n→ median_coverage", GREEN_ACC),
    ("Выравненные риды", "\u2265 80%", "Samtools stats\n→ reads_mapped_percent", ACCENT),
    ("Контаминация (TB-Mix)", "< 5%", "tb_mix.py — детекция\nсмешанных инфекций", ORANGE_ACC),
]):
    x = .5 + i * 4.0
    _card(sl, x, 3.05, 3.5, 1.2, clr)
    _tb(sl, x+.12, 3.1, 3.2, .18, lbl.upper(), Pt(9), LIGHT_GRAY, bold=True, align=PP_ALIGN.CENTER)
    _tb(sl, x+.12, 3.35, 3.2, .45, val, Pt(26), clr, bold=True, align=PP_ALIGN.CENTER)
    _tb(sl, x+.12, 3.85, 3.2, .35, desc, Pt(10), TEXT_GRAY, align=PP_ALIGN.CENTER)

_tb(sl, .7, 4.5, 10, .25, "Как работает TB-Mix:", Pt(15), TEXT_GRAY, bold=True)
_bullets(sl, .7, 4.8, 11, 2.0, [
    "Анализирует BAM-файл на гетерогенность аллельных частот в информативных SNP-позициях",
    "Порог смешанной инфекции: минорный аллель 5%–95% (--mix-low 0.05, --mix-high 0.95)",
    "Параметры фильтрации ридов: MQ \u2265 30, BQ \u2265 20, мин. частота аллеля \u2265 4%",
    "Образец со смешанной инфекцией \u2265 5% → исключается из дальнейшего анализа",
], Pt(13))

_num(sl, 6)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7 — ГЕНОТИПИРОВАНИЕ И РЕЗУЛЬТАТЫ (бывший слайд 4)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _title(sl, "Генотипирование и результаты")

_tb(sl, .7, 1.15, 5.5, .3, "6 параллельных модулей генотипирования", Pt(16), ACCENT, bold=True)
_tbl(sl, .7, 1.5, 5.8, [
    ["Модуль",      "Результат"],
    ["TBLG",        "Филогенетическая линия (L1–L7, подлинии)"],
    ["SpoTyping",   "Споликотип (43-спейсерный бинарный код)"],
    ["TB-Profiler", "Лекарственная устойчивость (RIF, INH, EMB...)"],
    ["RD-анализ",   "Регионы делеций (RD105, RD239, RD750...)"],
    ["ISMapper",    "Позиции IS6110 (только paired-end)"],
    ["Mosdepth",    "Профиль глубины покрытия по геному"],
], [1.4, 4.4], Pt(12))

_tb(sl, .7, 4.1, 5.8, .55,
    "Variant calling: Freebayes (ploidy=1, AF \u2265 80%, cov \u2265 10X)\n"
    "Аннотация: SnpEff (эффект мутаций, аминокислотные замены)",
    Pt(12), TEXT_GRAY)

_tb(sl, 7.0, 1.15, 5.5, .3, "Выходные файлы", Pt(16), ACCENT, bold=True)
_tbl(sl, 7.0, 1.5, 5.8, [
    ["Файл",               "Описание"],
    ["FINAL_TABLE.xlsx",    "Сводная таблица всех образцов"],
    ["general.tsv",         "QC-метрики когорты"],
    ["dr.xlsx",             "Лекарственная устойчивость"],
    ["spotyping.total.tsv", "Споликотипы"],
    ["rd.tsv",              "Делеции (RD)"],
    ["filter.tbmix.tsv",    "TB-Mix + линии"],
    ["MultiQC HTML",        "Интерактивный QC-дашборд"],
], [2.2, 3.6], Pt(12))

_tb(sl, 7.0, 4.1, 5.8, .55,
    "При > 1 образце: ANN_TABLE — когортная матрица SNP\n"
    "(bcftools merge → SnpEff → SnpSift → TSV)",
    Pt(12), TEXT_GRAY)

_num(sl, 7)


# ══════════════════════════════════════════════════════════════
#  SLIDE 8 — ДЕПЛОЙ И ПРОИЗВОДИТЕЛЬНОСТЬ (бывший слайд 5)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _title(sl, "Деплой, ресурсы и производительность")

LX = .7
_tb(sl, LX, 1.15, 5, .3, "Варианты запуска", Pt(16), ACCENT, bold=True)

for i, (lbl, cmd, clr) in enumerate([
    ("Локально / HPC (Singularity)", "nextflow run main.nf -profile local -resume", GREEN_ACC),
    ("Сервер (Docker)",              "nextflow run main.nf -resume",                 ACCENT),
    ("Kubernetes",                   "nextflow run main.nf -profile k8s",            ORANGE_ACC),
]):
    y = 1.5 + i * .7
    _card(sl, LX, y, 5.3, .6, clr)
    _tb(sl, LX+.12, y+.04, 5, .16, lbl.upper(), Pt(9), LIGHT_GRAY, bold=True)
    _tb(sl, LX+.12, y+.25, 5, .28, cmd, Pt(12), TEXT_DARK)

_tb(sl, LX, 3.7, 5, .25, "Контейнеры", Pt(16), ACCENT, bold=True)
_bullets(sl, LX, 4.0, 5, .7, [
    "15 кастомных Docker-образов + 8 из Docker Hub",
    "Поддержка Singularity (.sif) и Docker",
], Pt(13))

# Правая колонка: производительность
RX = 6.5; RW = 6.3
_tb(sl, RX, 1.15, RW, .3, "Время на 1 образец (полный прогон, 16 обр.)", Pt(15), ACCENT, bold=True)

_tbl(sl, RX, 1.5, RW, [
    ["Этап",                  "Среднее",  "Макс.",   "Peak RAM"],
    ["1. FASTP (тримминг)",   "28 сек",   "54 сек",  "1.5 ГБ"],
    ["2. FastQC",             "11 сек",   "27 сек",  "1.6 ГБ"],
    ["3. BWA + Picard",       "1.0 мин",  "3.2 мин", "3.5 ГБ"],
    ["4. FILTER (3 модуля)",  "55 сек",   "1.7 мин", "692 МБ"],
    ["5. Freebayes + SnpEff", "51 сек",   "1.5 мин", "270 МБ"],
    ["6. Genotyping (6 мод.)", "57 сек",  "2.9 мин", "979 МБ"],
    ["Итого на 1 обр.",       "~5.5 мин", "",        "bottleneck: BWA+Picard"],
], [1.8, 1.0, 1.0, 1.8], Pt(11), .3)

CW = 2.95; CH = .78; CG = .12
cx1 = RX; cx2 = RX + CW + CG
cy1 = 4.3; cy2 = cy1 + CH + .1

perf = [
    (cx1, cy1, "16 образцов (полный)",     "24 мин 39 сек", "wall clock, 244 задачи",   GREEN_ACC),
    (cx2, cy1, "100 обр. (экстраполяция)", "~1.5–2.5 ч",    "12 параллельных задач",    ACCENT),
    (cx1, cy2, "Оперативная память",       "\u2265 8 ГБ",    "пик: BWA+Picard 3.5 ГБ",  ORANGE_ACC),
    (cx2, cy2, "Дисковое пространство",   "~1–2 ГБ/обр.",  "100 обр. \u2248 100–200 ГБ", TEXT_GRAY),
]
for cx, cy, lbl, val, desc, clr in perf:
    _card(sl, cx, cy, CW, CH, clr)
    _tb(sl, cx+.1, cy+.04, CW-.2, .15, lbl.upper(), Pt(8), LIGHT_GRAY, bold=True, align=PP_ALIGN.CENTER)
    _tb(sl, cx+.1, cy+.22, CW-.2, .28, val, Pt(18), clr, bold=True, align=PP_ALIGN.CENTER)
    _tb(sl, cx+.1, cy+.54, CW-.2, .18, desc, Pt(9), TEXT_GRAY, align=PP_ALIGN.CENTER)

_num(sl, 8)


# ══════════════════════════════════════════════════════════════
#  SLIDE 9 — РЕЗУЛЬТАТЫ: КОЛЛЕКЦИЯ (данные из слайда 12)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _title(sl, "Результаты: коллекция геномов")

# Левая колонка: линии
_tb(sl, .7, 1.2, 6, .3, "Распределение по линиям (n = 79 094)", Pt(18), ACCENT, bold=True)

lineages = [
    ("L4",       36610, 46.3, BAR_BLUE),
    ("L2",       22995, 29.1, RGBColor(0x5D, 0xA5, 0xDA)),
    ("L3",        6797,  8.6, RGBColor(0x60, 0xBD, 0x68)),
    ("L1",        6080,  7.7, RGBColor(0xFA, 0xB7, 0x63)),
    ("M. bovis",  4837,  6.1, RGBColor(0xF1, 0x7C, 0xB0)),
    ("Другие",    1775,  2.2, LIGHT_GRAY),
]

for i, (nm, cnt, pct, clr) in enumerate(lineages):
    y = 1.65 + i * .5
    _tb(sl, .7, y, 1.3, .3, nm, Pt(15), TEXT_DARK, bold=True)
    bw = _bar(sl, 2.1, y + .06, 4.5, pct / 50.0, .28, clr)
    _tb(sl, 2.1 + bw + .1, y, 2.5, .3, f"{cnt:,} ({pct}%)", Pt(13), TEXT_GRAY)

_tb(sl, .7, 4.8, 6, .5,
    "Доминирование современных линий L4 и L2\nсогласуется с данными мировой литературы",
    Pt(13), TEXT_GRAY, italic=True)

# Правая колонка: страны
_tb(sl, 7.2, 1.2, 5.5, .3, "Первые 10 стран происхождения", Pt(18), ACCENT, bold=True)

_tbl(sl, 7.2, 1.6, 5.5, [
    ["Страна",       "Изоляты",  "%"],
    ["Georgia",      "4 946",    "12,0"],
    ["USA",          "4 151",    "10,0"],
    ["Australia",    "3 598",    "8,7"],
    ["China",        "3 349",    "8,1"],
    ["Moldova",      "3 251",    "7,9"],
    ["Peru",         "1 854",    "4,5"],
    ["Gambia",       "1 837",    "4,4"],
    ["South Africa", "1 329",    "3,2"],
    ["Viet Nam",     "1 132",    "2,7"],
    ["Brazil",       "1 123",    "2,7"],
], [2.0, 1.5, 1.0], Pt(13), .33)

_tb(sl, 7.2, 5.5, 5.5, .5,
    "Распределение отражает как эпидемиологию,\n"
    "так и доступность секвенирования и депонирования данных",
    Pt(12), TEXT_GRAY, italic=True)

_num(sl, 9)


# ══════════════════════════════════════════════════════════════
#  SLIDE 10 — РЕЗУЛЬТАТЫ: ЛЕКАРСТВЕННАЯ УСТОЙЧИВОСТЬ
#  (данные из слайда 13: категории + все 18 препаратов)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _title(sl, "Результаты: лекарственная устойчивость")

# Категории
_tb(sl, .7, 1.15, 12, .3, "Классификация устойчивости (n = 63 258)", Pt(18), ACCENT, bold=True)

cats = [
    ("DS-TB",    43011, "68,0%", GREEN_ACC),
    ("MDR-TB",   12515, "19,8%", ORANGE_ACC),
    ("pre-XDR",   7189, "11,4%", RED_ACC),
    ("XDR-TB",     543, "0,9%",  DARK_RED),
]
for i, (nm, cnt, pct, clr) in enumerate(cats):
    x = .7 + i * 3.1
    _card(sl, x, 1.5, 2.8, .75, clr)
    _tb(sl, x+.08, 1.53, 2.6, .2, nm, Pt(14), clr, bold=True, align=PP_ALIGN.CENTER)
    _tb(sl, x+.08, 1.75, 2.6, .28, pct, Pt(22), clr, bold=True, align=PP_ALIGN.CENTER)
    _tb(sl, x+.08, 2.02, 2.6, .18, f"{cnt:,}", Pt(10), TEXT_GRAY, align=PP_ALIGN.CENTER)

_tb(sl, .7, 2.4, 12, .2,
    "Определения pre-XDR/XDR соответствуют обновлённой классификации ВОЗ (2021)",
    Pt(11), TEXT_GRAY, italic=True)

# Все 18 препаратов (две колонки по 9)
_tb(sl, .7, 2.75, 12, .3, "Устойчивость по 18 противотуберкулёзным препаратам", Pt(18), ACCENT, bold=True)

drugs = [
    ("Isoniazid",     25848, RED_ACC),
    ("Rifampicin",    21933, RED_ACC),
    ("Streptomycin",  21890, RED_ACC),
    ("Pyrazinamide",  18190, ORANGE_ACC),
    ("Ethambutol",    17683, ORANGE_ACC),
    ("Ethionamide",   13447, ORANGE_ACC),
    ("Moxifloxacin",   8512, BAR_BLUE),
    ("Levofloxacin",   8512, BAR_BLUE),
    ("Kanamycin",      6032, BAR_BLUE),
    ("Capreomycin",    3857, BAR_BLUE),
    ("Amikacin",       3649, BAR_BLUE),
    ("PAS",            1761, LIGHT_GRAY),
    ("Bedaquiline",     692, LIGHT_GRAY),
    ("Clofazimine",     676, LIGHT_GRAY),
    ("Delamanid",       290, LIGHT_GRAY),
    ("Linezolid",       242, LIGHT_GRAY),
    ("Pretomanid",      241, LIGHT_GRAY),
    ("Cycloserine",     231, LIGHT_GRAY),
]

max_n = 25848  # максимум (isoniazid)
for col in range(2):
    x_name = .7 + col * 6.3
    x_cnt  = x_name + 1.7
    x_bar  = x_cnt + .8
    bar_max = 2.3
    x_pct  = x_bar + bar_max + .1

    for i in range(9):
        idx = col * 9 + i
        if idx >= len(drugs):
            break
        nm, cnt, clr = drugs[idx]
        pct = cnt / 63258 * 100
        y = 3.15 + i * .43

        _tb(sl, x_name, y, 1.65, .28, nm, Pt(12), TEXT_DARK)
        _tb(sl, x_cnt, y, .75, .28, f"{cnt:,}", Pt(12), TEXT_DARK, bold=True, align=PP_ALIGN.RIGHT)
        _bar(sl, x_bar, y + .05, bar_max, cnt / max_n, .2, clr)
        _tb(sl, x_pct, y, .7, .28, f"{pct:.1f}%", Pt(10), TEXT_GRAY)

_num(sl, 10)


# ══════════════════════════════════════════════════════════════
#  SLIDE 11 — ЗАКЛЮЧЕНИЕ
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _title(sl, "Заключение")

_nbullets(sl, .7, 1.25, 11.5, 4.8, [
    "Разработан Nextflow-пайплайн TB-Lite (8 этапов, 26 процессов)\n"
    "для комплексного анализа данных WGS M. tuberculosis",
    "Реализована трёхуровневая фильтрация (покрытие \u2265 30X,\n"
    "выравнивание \u2265 80%, контаминация < 5%)\n"
    "и модуль детекции смешанных инфекций TB-Mix",
    "Пайплайн объединяет типирование (TBLG, SpoTyping),\n"
    "оценку устойчивости (TB-Profiler), анализ делеций (RD)\n"
    "и локализацию IS6110 (ISMapper)",
    "Апробация на >80 000 геномов NCBI:\n"
    "стабильная работа, 16 образцов за ~25 мин",
    "Доминирование линий L4 (46,3%) и L2 (29,1%);\n"
    "устойчивость проанализирована по 18 препаратам\n"
    "для 63 258 изолятов",
    "Пайплайн пригоден для первичной обработки\n"
    "новых изолятов и повторного анализа публичных\n"
    "коллекций — стандартизация данных для геномных\n"
    "и эпидемиологических исследований",
], Pt(18), Pt(12))

_tb(sl, .7, 6.0, 11.5, .5, "Спасибо за внимание!",
    Pt(28), DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
_tb(sl, .7, 6.6, 11.5, .3, "GitHub: github.com/kiselevs/tb-lite",
    Pt(14), ACCENT, align=PP_ALIGN.CENTER)
_num(sl, 11)


# ─── Save ───
out = "/home/zerg/git/tb-lite/TB-Lite_Conference.pptx"
prs.save(out)
print(f"Saved: {out}")
